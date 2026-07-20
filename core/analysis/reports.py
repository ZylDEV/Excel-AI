"""PDF and PowerPoint report generator from dashboard data."""

import logging
import os
import warnings
from datetime import datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

REPORTS_DIR = Path(__file__).resolve().parent.parent.parent / "backend" / "reports"


class ReportGenerator:
    """Generate PDF and PPTX reports from dashboard configuration data."""

    @staticmethod
    def _ensure_reports_dir() -> Path:
        """Create the reports output directory if it does not exist."""
        REPORTS_DIR.mkdir(parents=True, exist_ok=True)
        return REPORTS_DIR

    @staticmethod
    def _sanitize_filename(name: str) -> str:
        """Remove characters that are problematic in filenames."""
        return "".join(c if c.isalnum() or c in " _.-" else "_" for c in name).strip()

    @classmethod
    def generate_report_data(cls, dashboard_data: dict) -> dict[str, Any]:
        """Prepare and normalise dashboard data for report generation.

        Parameters
        ----------
        dashboard_data : dict
            The output of ``DashboardGenerator.generate()``.

        Returns
        -------
        dict with cleaned/normalised keys: kpis, charts, insights, summary.
        """
        kpis = dashboard_data.get("kpis", [])
        charts = dashboard_data.get("charts", [])
        insights = dashboard_data.get("insights", [])
        summary = dashboard_data.get("summary", {})

        # Ensure KPIs are flat dicts with sensible defaults
        cleaned_kpis = []
        for kpi in kpis:
            cleaned_kpis.append({
                "label": kpi.get("label", ""),
                "value": kpi.get("value", ""),
                "change_pct": kpi.get("change_pct", 0.0),
                "icon": kpi.get("icon", ""),
                "color": kpi.get("color", "#333"),
                "prefix": kpi.get("prefix", ""),
                "suffix": kpi.get("suffix", ""),
            })

        return {
            "kpis": cleaned_kpis,
            "charts": charts,
            "insights": insights if isinstance(insights, list) else [str(insights)],
            "summary": summary if isinstance(summary, dict) else {},
        }

    @classmethod
    def generate_pdf(cls, dashboard_data: dict, output_path: str | None = None) -> str:
        """Generate a PDF report from dashboard data.

        Uses fpdf2 to create a multi-page report with:
        - Cover page
        - Executive summary
        - KPI table
        - Chart representations (table-based)
        - Insights list

        Parameters
        ----------
        dashboard_data : dict
            Dashboard configuration data.
        output_path : str | None
            Full path for the output file. If None, a timestamped name is
            used inside the report directory.

        Returns
        -------
        str — the absolute path to the generated PDF file.
        """
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", UserWarning)
            try:
                from fpdf import FPDF

                cls._ensure_reports_dir()

                if output_path is None:
                    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_path = str(REPORTS_DIR / f"report_{ts}.pdf")

                data = cls.generate_report_data(dashboard_data)

                pdf = FPDF(orientation="P", unit="mm", format="A4")
                pdf.set_auto_page_break(auto=True, margin=20)

                # ── Cover Page ────────────────────────────────────────
                pdf.add_page()
                pdf.set_font("Helvetica", "B", 28)
                pdf.set_text_color(33, 150, 243)
                pdf.cell(0, 40, "Excel AI", align="C", new_x="LMARGIN", new_y="NEXT")
                pdf.set_font("Helvetica", "", 18)
                pdf.set_text_color(80, 80, 80)
                pdf.cell(0, 15, "Laporan Dashboard", align="C", new_x="LMARGIN", new_y="NEXT")
                pdf.ln(10)
                pdf.set_font("Helvetica", "", 12)
                pdf.set_text_color(120, 120, 120)
                pdf.cell(0, 10, f"Generated: {datetime.now().strftime('%d %B %Y %H:%M')}", align="C", new_x="LMARGIN", new_y="NEXT")

                # ── Executive Summary ─────────────────────────────────
                pdf.add_page()
                pdf.set_font("Helvetica", "B", 18)
                pdf.set_text_color(33, 33, 33)
                pdf.cell(0, 15, "Ringkasan Eksekutif", new_x="LMARGIN", new_y="NEXT")
                pdf.ln(5)
                pdf.set_font("Helvetica", "", 11)
                pdf.set_text_color(60, 60, 60)

                summary = data.get("summary", {})
                summary_lines = [
                    f"Total Baris: {summary.get('row_count', 'N/A')}",
                    f"Total Kolom: {summary.get('column_count', 'N/A')}",
                    f"Total Revenue: {summary.get('total_revenue', 'N/A')}",
                    f"Total Profit: {summary.get('total_profit', 'N/A')}",
                    f"Growth: {summary.get('growth', 0)}%",
                ]
                for line in summary_lines:
                    pdf.cell(0, 8, f"  - {line}", new_x="LMARGIN", new_y="NEXT")
                pdf.ln(10)

                # ── KPI Table ─────────────────────────────────────────
                pdf.set_font("Helvetica", "B", 18)
                pdf.set_text_color(33, 33, 33)
                pdf.cell(0, 15, "Indikator Kinerja Utama (KPI)", new_x="LMARGIN", new_y="NEXT")
                pdf.ln(5)

                kpis = data.get("kpis", [])
                if kpis:
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_fill_color(33, 150, 243)
                    pdf.set_text_color(255, 255, 255)
                    col_w = [80, 50, 50]
                    headers_table = ["Metrik", "Nilai", "Perubahan (%)"]
                    for i, h in enumerate(headers_table):
                        pdf.cell(col_w[i], 10, h, border=1, fill=True, align="C")
                    pdf.ln()

                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(60, 60, 60)
                    for kpi in kpis[:10]:
                        label = str(kpi.get("label", ""))[:35]
                        value = str(kpi.get("value", ""))
                        change = f"{kpi.get('change_pct', 0):+.2f}%"
                        pdf.cell(col_w[0], 8, label, border=1)
                        pdf.cell(col_w[1], 8, value, border=1, align="C")
                        pdf.cell(col_w[2], 8, change, border=1, align="C")
                        pdf.ln()
                else:
                    pdf.set_font("Helvetica", "", 11)
                    pdf.cell(0, 8, "Tidak ada KPI tersedia.", new_x="LMARGIN", new_y="NEXT")

                pdf.ln(10)

                # ── Insights ──────────────────────────────────────────
                pdf.set_font("Helvetica", "B", 18)
                pdf.set_text_color(33, 33, 33)
                pdf.cell(0, 15, "Wawasan Utama", new_x="LMARGIN", new_y="NEXT")
                pdf.ln(5)
                pdf.set_font("Helvetica", "", 11)
                pdf.set_text_color(60, 60, 60)

                insights = data.get("insights", [])
                if insights:
                    for i, insight in enumerate(insights, 1):
                        pdf.multi_cell(0, 8, f"{i}. {insight}")
                        pdf.ln(2)
                else:
                    pdf.cell(0, 8, "Tidak ada wawasan tersedia.", new_x="LMARGIN", new_y="NEXT")

                pdf.output(output_path)
                logger.info(f"PDF report generated: {output_path}")
                return str(Path(output_path).resolve())

            except Exception as e:
                logger.exception("Gagal membuat PDF report")
                raise RuntimeError(f"Gagal membuat PDF: {e}") from e

    @classmethod
    def generate_ppt(cls, dashboard_data: dict, output_path: str | None = None) -> str:
        """Generate a PowerPoint report from dashboard data.

        Uses python-pptx to create slides with:
        - Title slide
        - Executive summary slide
        - KPI summary slide(s)
        - Chart slides (one per chart)
        - Insights slide

        Parameters
        ----------
        dashboard_data : dict
            Dashboard configuration data.
        output_path : str | None
            Full path for the output file.

        Returns
        -------
        str — the absolute path to the generated PPTX file.
        """
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", UserWarning)
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN

                cls._ensure_reports_dir()

                if output_path is None:
                    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                    output_path = str(REPORTS_DIR / f"report_{ts}.pptx")

                data = cls.generate_report_data(dashboard_data)

                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)

                BLUE = RGBColor(33, 150, 243)
                DARK = RGBColor(33, 33, 33)
                GRAY = RGBColor(100, 100, 100)
                WHITE = RGBColor(255, 255, 255)

                def _add_textbox(slide, left, top, width, height, text, font_size=18,
                                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(font_size)
                    p.font.bold = bold
                    p.font.color.rgb = color
                    p.alignment = alignment
                    return txBox

                def _add_bg(slide, color):
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = color

                # ── Slide 1: Title ────────────────────────────────────
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
                _add_bg(slide, BLUE)
                _add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                             "Excel AI", font_size=44, bold=True, color=WHITE,
                             alignment=PP_ALIGN.CENTER)
                _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                             "Laporan Dashboard", font_size=28, color=WHITE,
                             alignment=PP_ALIGN.CENTER)
                _add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(1),
                             datetime.now().strftime("%d %B %Y %H:%M"), font_size=16,
                             color=WHITE, alignment=PP_ALIGN.CENTER)

                # ── Slide 2: Executive Summary ────────────────────────
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                             "Ringkasan Eksekutif", font_size=32, bold=True, color=BLUE)

                summary = data.get("summary", {})
                summary_items = [
                    f"Total Baris: {summary.get('row_count', 'N/A')}",
                    f"Total Kolom: {summary.get('column_count', 'N/A')}",
                    f"Total Revenue: {summary.get('total_revenue', 'N/A')}",
                    f"Total Profit: {summary.get('total_profit', 'N/A')}",
                    f"Growth: {summary.get('growth', 0)}%",
                ]
                y_pos = 1.5
                for item in summary_items:
                    _add_textbox(slide, Inches(1), Inches(y_pos), Inches(11), Inches(0.5),
                                 f"  - {item}", font_size=18, color=DARK)
                    y_pos += 0.6

                # ── Slide 3: KPIs ─────────────────────────────────────
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                             "Indikator Kinerja Utama (KPI)", font_size=28, bold=True,
                             color=BLUE)

                kpis = data.get("kpis", [])
                y_pos = 1.5
                for kpi in kpis[:12]:
                    label = str(kpi.get("label", ""))
                    value = str(kpi.get("value", ""))
                    change = f"{kpi.get('change_pct', 0):+.2f}%"
                    _add_textbox(slide, Inches(0.5), Inches(y_pos), Inches(4), Inches(0.5),
                                 label, font_size=14, bold=True, color=DARK)
                    _add_textbox(slide, Inches(4.5), Inches(y_pos), Inches(3), Inches(0.5),
                                 value, font_size=14, color=DARK)
                    _add_textbox(slide, Inches(7.5), Inches(y_pos), Inches(3), Inches(0.5),
                                 change, font_size=14,
                                 color=RGBColor(76, 175, 80) if kpi.get("change_pct", 0) >= 0
                                 else RGBColor(244, 67, 54))
                    y_pos += 0.4

                # ── Slide 4+: Charts ──────────────────────────────────
                charts = data.get("charts", [])
                for chart in charts:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    title = chart.get("title", "Chart")
                    chart_type = chart.get("type", "bar")
                    labels = chart.get("labels", [])
                    datasets = chart.get("datasets", [])

                    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                                 f"Chart: {title}", font_size=28, bold=True, color=BLUE)

                    # Description of the chart
                    chart_desc = (
                        f"Tipe: {chart_type.upper()} | "
                        f"Jumlah kategori: {len(labels)} | "
                        f"Jumlah series: {len(datasets)}"
                    )
                    _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                                 chart_desc, font_size=14, color=GRAY)

                    # Show top data points
                    _add_textbox(slide, Inches(0.5), Inches(2), Inches(12), Inches(0.5),
                                 "Data Points:", font_size=16, bold=True, color=DARK)

                    y_pos = 2.8
                    for i, label in enumerate(labels[:10]):
                        val_str = "; ".join(
                            f"{d['label']}: {d['data'][i]}" for d in datasets if i < len(d.get("data", []))
                        )
                        _add_textbox(slide, Inches(0.8), Inches(y_pos), Inches(11), Inches(0.4),
                                     f"  {i + 1}. {label} - {val_str}",
                                     font_size=12, color=DARK)
                        y_pos += 0.35

                # ── Final slide: Insights ─────────────────────────────
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                             "Wawasan Utama", font_size=28, bold=True, color=BLUE)

                insights = data.get("insights", [])
                y_pos = 1.5
                if insights:
                    for i, insight in enumerate(insights, 1):
                        _add_textbox(slide, Inches(0.8), Inches(y_pos), Inches(11), Inches(0.6),
                                     f"{i}. {insight}", font_size=16, color=DARK)
                        y_pos += 0.6
                else:
                    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                                 "Tidak ada wawasan tersedia.", font_size=16, color=GRAY)

                prs.save(output_path)
                logger.info(f"PPTX report generated: {output_path}")
                return str(Path(output_path).resolve())

            except Exception as e:
                logger.exception("Gagal membuat PPT report")
                raise RuntimeError(f"Gagal membuat PPT: {e}") from e
